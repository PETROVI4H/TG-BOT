import telebot
import os
import requests
import re
from telebot import types
from io import BytesIO
import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from fuzzywuzzy import fuzz
from pptx import Presentation

bot = telebot.TeleBot('8131849234:AAEAduQ2nGxxaAgWRsW35nf9Lq1YoavNZfk')

bot.remove_webhook()

uploaded_file_path = None
df = None
all_matching_rows = []

def sanitize_filename(file_url):
    return re.sub(r'[<>:"/\\|?*#]', '_', file_url)

@bot.message_handler(commands=['start'])
def start(message):
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton('Перейти в Omni', url='https://omni.top-academy.ru/login/index#/'))

    welcome_message = (
        f'Здравствуйте! {message.from_user.first_name} {message.from_user.last_name}\n\n'
        'Я помогу вам работать с документами. Для начала, пожалуйста, загрузите файл или отправьте ссылку на документ.'
        '\nПосле загрузки вы сможете выбрать команду для поиска данных в документе.'
    )

    bot.send_message(message.chat.id, welcome_message, reply_markup=markup)

@bot.message_handler(commands=['help'])
def about(message):
    help_message = (
        "<b><u>Информация о боте:</u></b>\n\n"
        "Этот бот поможет вам работать с различными документами, включая Excel, Google Docs и другие. Вот список доступных команд:\n\n"
        "<b>/start</b> - Стартовая команда, которая приветствует вас и дает инструкции по загрузке документа.\n"
        "<b>/web</b> - Открывает официальный сайт преподавателя OMNI.\n"
        "<b>/help</b> - Выводит информацию о доступных функциях бота.\n"
        "<b>/more</b> - Позволяет получить дополнительные результаты поиска.\n"
        "<b>/url</b> - Позволяет отправить ссылку на файл (Google Docs, PDF, Word и другие), который бот загрузит и проанализирует.\n\n"
        "<b>После загрузки файла:</b>\n"
        "Вы сможете использовать команды справа внизу для поиска текста.\n"
    )
    bot.send_message(message.chat.id, help_message, parse_mode='html')

@bot.message_handler(commands=['url'])
def ask_for_url(message):
    bot.send_message(message.chat.id, "Пожалуйста, отправьте ссылку на файл (Google Docs, PDF, Word и другие), чтобы я мог его загрузить и проанализировать.")

@bot.message_handler(func=lambda message: message.text.startswith('http'))
def handle_url(message):
    global uploaded_file_path, df, all_matching_rows

    file_url = message.text.strip()

    try:
        response = requests.get(file_url)

        if response.status_code == 200:
            content_type = response.headers.get('Content-Type', '').lower()

            if 'excel' in content_type:
                user_id = message.from_user.id
                user_dir = f'./user_files/{user_id}'

                os.makedirs(user_dir, exist_ok=True)

                file_name = sanitize_filename(file_url.split("/")[-1])
                file_path = os.path.join(user_dir, file_name)

                with open(file_path, 'wb') as f:
                    f.write(response.content)

                try:
                    df = pd.read_excel(file_path, engine="openpyxl")
                    uploaded_file_path = file_path
                    all_matching_rows = []
                    bot.send_message(message.chat.id, f"Файл Excel загружен: {file_name}. Пример данных:\n{df.head()}")
                except Exception as e:
                    bot.send_message(message.chat.id, f"Ошибка при обработке Excel файла: {e}")
                    return

            elif 'pdf' in content_type:
                bot.send_message(message.chat.id, "Я сейчас не могу анализировать PDF файлы.")
            elif 'msword' in content_type or 'word' in content_type:
                doc = Document(BytesIO(response.content))
                full_text = '\n'.join([para.text for para in doc.paragraphs])
                bot.send_message(message.chat.id, f"Текст из Word документа:\n{full_text[:1000]}")
            elif 'csv' in content_type:
                user_id = message.from_user.id
                user_dir = f'./user_files/{user_id}'

                os.makedirs(user_dir, exist_ok=True)

                file_name = sanitize_filename(file_url.split("/")[-1])
                file_path = os.path.join(user_dir, file_name)

                with open(file_path, 'wb') as f:
                    f.write(response.content)

                try:
                    df = pd.read_csv(file_path)
                    uploaded_file_path = file_path
                    all_matching_rows = []
                    bot.send_message(message.chat.id, f"CSV файл загружен: {file_name}. Пример данных:\n{df.head()}")
                except Exception as e:
                    bot.send_message(message.chat.id, f"Ошибка при обработке CSV файла: {e}")
                    return
            elif 'vnd.ms-powerpoint' in content_type or 'presentation' in content_type:
                user_id = message.from_user.id
                user_dir = f'./user_files/{user_id}'

                os.makedirs(user_dir, exist_ok=True)

                file_name = sanitize_filename(file_url.split("/")[-1])
                file_path = os.path.join(user_dir, file_name)

                with open(file_path, 'wb') as f:
                    f.write(response.content)

                try:
                    prs = Presentation(file_path)
                    all_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                all_text.append(shape.text)
                    bot.send_message(message.chat.id, "Текст из PowerPoint презентации:\n" + "\n".join(all_text[:1000]))
                except Exception as e:
                    bot.send_message(message.chat.id, f"Ошибка при обработке PowerPoint файла: {e}")
                    return
            else:
                bot.send_message(message.chat.id, "Файл имеет неподдерживаемый формат или не может быть проанализирован.")
        else:
            bot.send_message(message.chat.id, "Не удалось загрузить файл по предоставленной ссылке. Пожалуйста, проверьте ссылку.")
    except Exception as e:
        bot.send_message(message.chat.id, f"Ошибка при обработке файла: {e}")
        return

@bot.message_handler(content_types=['document'])
def handle_document(message):
    global uploaded_file_path, df, all_matching_rows

    user_id = message.from_user.id
    user_dir = f'./user_files/{user_id}'

    os.makedirs(user_dir, exist_ok=True)

    file_info = bot.get_file(message.document.file_id)
    file = bot.download_file(file_info.file_path)

    file_name = message.document.file_name
    file_path = os.path.join(user_dir, file_name)

    with open(file_path, 'wb') as f:
        f.write(file)

    try:
        if file_name.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(file_path, engine="openpyxl")
            uploaded_file_path = file_path
            all_matching_rows = []
            bot.send_message(message.chat.id, f"Файл загружен и сохранен как {file_name}. Пример данных:\n{df.head()}")
            markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
            markup.add(types.KeyboardButton("Проверить наличие текста в файле"))
            markup.add(types.KeyboardButton("Проверить правильность текста"))
            bot.send_message(message.chat.id, "Выберите один из вариантов:", reply_markup=markup)
        elif file_name.endswith('.csv'):
            df = pd.read_csv(file_path)
            uploaded_file_path = file_path
            all_matching_rows = []
            bot.send_message(message.chat.id, f"CSV файл загружен и сохранен как {file_name}. Пример данных:\n{df.head()}")
            markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
            markup.add(types.KeyboardButton("Проверить наличие текста в файле"))
            markup.add(types.KeyboardButton("Проверить правильность текста"))
            bot.send_message(message.chat.id, "Выберите один из вариантов:", reply_markup=markup)
        elif file_name.endswith('.pdf'):
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                full_text = ""
                for page in reader.pages:
                    full_text += page.extract_text()
            bot.send_message(message.chat.id, f"Текст из PDF файла:\n{full_text[:1000]}")
        elif file_name.endswith('.pptx'):
            prs = Presentation(file_path)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            bot.send_message(message.chat.id, "Текст из PowerPoint файла:\n" + "\n".join(all_text[:1000]))
        elif file_name.endswith('.docx'):
            doc = Document(file_path)
            full_text = '\n'.join([para.text for para in doc.paragraphs])
            bot.send_message(message.chat.id, f"Текст из Word документа:\n{full_text[:1000]}")
        else:
            bot.send_message(message.chat.id, "Поддерживаются только Excel файлы, CSV, текстовые файлы и PowerPoint.")
    except Exception as e:
        bot.send_message(message.chat.id, f"Ошибка при обработке файла: {e}")

@bot.message_handler(func=lambda message: message.text == "Проверить наличие текста в файле")
def compare_text(message):
    global df, all_matching_rows

    if df is None:
        bot.send_message(message.chat.id, "Сначала загрузите файл, чтобы я мог выполнить поиск.")
        return

    bot.send_message(message.chat.id, "Введите текст для сравнения с данными из документа:")

    bot.register_next_step_handler(message, process_comparison)

@bot.message_handler(func=lambda message: message.text == "Проверить правильность текста")
def check_text_correctness(message):
    if df is None:
        bot.send_message(message.chat.id, "Сначала загрузите файл, чтобы я мог проверить правильность текста.")
        return

    bot.send_message(message.chat.id, "Введите текст для проверки на наличие ошибок:")

    bot.register_next_step_handler(message, process_text_correction)

def process_comparison(message):
    global df, all_matching_rows

    user_text = message.text.strip()
    matching_rows = []

    if isinstance(df, pd.DataFrame):
        for column in df.columns:
            column_data = df[column].astype(str)

            matching_column_rows = df[column_data.str.contains(user_text, na=False, case=False, regex=True)]

            if not matching_column_rows.empty:
                for index, row in matching_column_rows.iterrows():
                    context = f"Столбец: {column}\n"
                    context += f"Строка {index}: {row[column]}\n"
                    matching_rows.append(context)

    if matching_rows:
        all_matching_rows = matching_rows
        bot.send_message(message.chat.id, "Текст найден! Вот примеры:\n")
        for row in matching_rows[:5]:
            bot.send_message(message.chat.id, row)
        if len(matching_rows) > 5:
            bot.send_message(message.chat.id, "... и другие совпадения. Вы можете запросить больше с помощью команды /more.")
    else:
        bot.send_message(message.chat.id, "Текст не найден.")

def process_text_correction(message):
    global df

    user_text = message.text.strip()
    best_match = None
    best_score = 0

    if isinstance(df, pd.DataFrame):
        for column in df.columns:
            column_data = df[column].astype(str)

            for value in column_data:
                score = fuzz.ratio(user_text.lower(), value.lower())
                if score > best_score:
                    best_score = score
                    best_match = value

    if best_score >= 80:
        bot.send_message(message.chat.id, f"Текст совпадает с файлом! Пример: {best_match}")
    else:
        bot.send_message(message.chat.id, "Текст написан неверно. Вот строка, которая максимально похожа:\n" + best_match)

@bot.message_handler(commands=['more'])
def more_results(message):
    global all_matching_rows

    if not all_matching_rows:
        bot.send_message(message.chat.id, "Сначала выполните команду сравнения справа внизу, чтобы найти совпадения.")
        return

    bot.send_message(message.chat.id, "Вот все совпадения:")

    for row in all_matching_rows:
        bot.send_message(message.chat.id, row)

    bot.send_message(message.chat.id, "Это все совпадения. Если вам нужно что-то еще, напишите.")

bot.polling(none_stop=True)